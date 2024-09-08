import pandas as pd
import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# Load the CSV file
df = pd.read_csv('Master Tracker 2024 krk(Movement).csv')

# Normalize column names by stripping whitespace
df.columns = df.columns.str.strip()

# Parse 'Sending Date' and filter data to include records up to August 2024
df['Sending Date'] = pd.to_datetime(df['Sending Date'], errors='coerce')
df = df[df['Sending Date'] <= '2024-08-31']

# Filter records for outbound movements
outbound_movements = df[df['Type Of Movement'] == 'Outbound']

# Count of orders by warehouse
warehouse_counts = df['Warehouse'].value_counts()
top_n_warehouses = warehouse_counts.nlargest(10)

# Create bar chart for top warehouses
plt.figure(figsize=(10, 6))
top_n_warehouses.plot(kind='bar', color='lightgreen')
plt.title('Top 10 Warehouses by Order Count (Up to August 2024)')
plt.xlabel('Warehouse')
plt.ylabel('Count')
plt.xticks(rotation=45)
plt.tight_layout()
plt.savefig('top_warehouses_bar_chart.png')
plt.close()

# Create pie chart for distribution of orders by top warehouses
plt.figure(figsize=(8, 8))
top_n_warehouses.plot(kind='pie', autopct='%1.1f%%', startangle=90, colors=plt.cm.Paired.colors)
plt.title('Distribution of Orders by Top 10 Warehouses (Up to August 2024)')
plt.ylabel('')
plt.tight_layout()
plt.savefig('warehouse_distribution_pie_chart.png')
plt.close()

# Count of orders by status
status_counts = df['SR Statues'].value_counts()

# Create bar chart for orders by status
plt.figure(figsize=(10, 6))
status_counts.plot(kind='bar', color='skyblue')
plt.title('Count of Orders by Status (Up to August 2024)')
plt.xlabel('Status')
plt.ylabel('Count')
plt.xticks(rotation=45)
plt.tight_layout()
plt.savefig('orders_by_status_bar_chart.png')
plt.close()

# Group outbound movements by 'Sending Date'
outbound_counts = outbound_movements.groupby(outbound_movements['Sending Date'].dt.to_period('M')).size().reset_index(name='Count')
outbound_counts['Sending Date'] = outbound_counts['Sending Date'].dt.to_timestamp()

# Create line chart for outbound movements over time
plt.figure(figsize=(12, 6))
plt.plot(outbound_counts['Sending Date'], outbound_counts['Count'], marker='o', color='orange')
plt.title('Outbound Movements Over Time (Up to August 2024)')
plt.xlabel('Date')
plt.ylabel('Count of Outbound Cases')
plt.xticks(rotation=45)
plt.tight_layout()
plt.savefig('outbound_movements_over_time.png')
plt.close()

# Convert 'Qty's cases' to numeric and drop NaN values
outbound_movements.loc[:, "Qty's cases"] = pd.to_numeric(outbound_movements["Qty's cases"], errors='coerce')
outbound_movements = outbound_movements.dropna(subset=["Qty's cases"])

# Create histogram for quantity distribution
plt.figure(figsize=(10, 6))
outbound_movements["Qty's cases"].plot(kind='hist', bins=30, color='lightcoral', edgecolor='black')
plt.title('Distribution of Quantities for Outbound Movements (Up to August 2024)')
plt.xlabel('Quantity of Cases')
plt.ylabel('Frequency')
plt.tight_layout()
plt.savefig('quantity_distribution_histogram.png')
plt.close()

# Client Movement Analysis
movement_counts = df.groupby(['Clients', 'Type Of Movement']).size().unstack(fill_value=0)
top_n = 10
top_clients = movement_counts.sum(axis=1).nlargest(top_n).index
filtered_movement_counts = movement_counts.loc[top_clients]

# Create bar chart for inbound and outbound movements by top clients
plt.figure(figsize=(15, 8))
filtered_movement_counts.plot(kind='bar', stacked=True)
plt.title(f'Top {top_n} Clients: Inbound and Outbound Movements (Up to August 2024)')
plt.xlabel('Clients')
plt.ylabel('Number of Movements')
plt.xticks(rotation=45, ha='right')
plt.legend(title='Type Of Movement')
plt.tight_layout()
plt.savefig('top_clients_movements_bar_chart.png')
plt.close()

# Create pie chart for total inbound vs outbound
inbound_counts = filtered_movement_counts['Inbound'].sum()
outbound_counts = filtered_movement_counts['Outbound'].sum()
plt.figure(figsize=(10, 6))
plt.pie([inbound_counts, outbound_counts], labels=['Total Inbound', 'Total Outbound'], autopct='%1.1f%%', startangle=90)
plt.title('Total Inbound vs Outbound Movements (Top Clients, Up to August 2024)')
plt.tight_layout()
plt.savefig('inbound_vs_outbound_pie_chart.png')
plt.close()

# Create PowerPoint presentation
prs = Presentation()

# Function to add a slide with a title and content
def add_slide(title, content, image_path=None):
    slide_layout = prs.slide_layouts[1]  # Use the title and content layout
    slide = prs.slides.add_slide(slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]

    title_placeholder.text = title
    content_placeholder.text = content

    if image_path:
        left = Inches(1)
        top = Inches(2.5)
        slide.shapes.add_picture(image_path, left, top, width=Inches(8))

# Title Slide
slide_layout = prs.slide_layouts[0]  # Use the title layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Analysis of Outbound Movements and Client Activity"
subtitle.text = "Insights from the Master Tracker 2024 Data\n[Your Name]\n[Date]"

# Add slides with content and charts
add_slide("Introduction", "Objective: To analyze outbound movements and client activities using the Master Tracker data.\nData Source: Master Tracker 2024 krk(Movement).csv")
add_slide("Data Cleaning and Preparation", "Steps Taken:\n- Loaded the CSV and normalized column names.\n- Filtered data to include records up to August 2024.")
add_slide("Outbound Movements Overview", "Total Records Analyzed: [Total Count]\nKey Metrics:\n- Number of outbound movements.\n- Breakdown by status and warehouse.")
add_slide("Top 10 Warehouses by Order Count", "Bar Chart: Top 10 Warehouses by Order Count (Up to August 2024)", 'top_warehouses_bar_chart.png')
add_slide("Distribution of Orders by Top Warehouses", "Pie Chart: Distribution of Orders by Top 10 Warehouses (Up to August 2024)", 'warehouse_distribution_pie_chart.png')
add_slide("Orders by Status", "Bar Chart: Count of Orders by Status (Up to August 2024)", 'orders_by_status_bar_chart.png')
add_slide("Outbound Movements Over Time", "Line Chart: Outbound Movements Over Time (Up to August 2024)", 'outbound_movements_over_time.png')
add_slide("Quantity Distribution for Outbound Movements", "Histogram: Distribution of Quantities for Outbound Movements (Up to August 2024)", 'quantity_distribution_histogram.png')
add_slide("Client Movement Analysis", "Top Clients Overview: Total inbound and outbound movements.")
add_slide("Inbound and Outbound Movements by Top Clients", "Bar Chart: Top 10 Clients: Inbound and Outbound Movements (Up to August 2024)", 'top_clients_movements_bar_chart.png')
add_slide("Total Inbound vs Outbound", "Pie Chart: Total Inbound vs Outbound Movements (Top Clients, Up to August 2024)", 'inbound_vs_outbound_pie_chart.png')

# Suggestions for Further Analysis
add_slide("Suggestions for Further Analysis", "- Implement advanced forecasting methods (e.g., ARIMA, exponential smoothing).\n- Analyze seasonal trends and their impact on outbound movements.\n- Explore correlations between order status and client behavior.")
# Potential Improvements
add_slide("Potential Improvements", "- Automate data cleaning and visualization processes.\n- Integrate additional data sources for comprehensive analysis (e.g., client feedback, shipping delays).\n- Develop a dashboard for real-time monitoring of outbound movements.")

# Conclusion
add_slide("Conclusion", "Summary of key findings.\nImportance of continuous monitoring and analysis for operational efficiency.")
# Questions
add_slide("Questions", "Open the floor for any questions or discussions.")

# Save the presentation
prs.save('Outbound_Movements_Analysis_Presentation.pptx')